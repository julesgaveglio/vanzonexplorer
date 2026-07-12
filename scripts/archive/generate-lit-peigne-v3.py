#!/usr/bin/env python3
"""
Guide VBA Construction — Lit Peigne v3
100% formes natives PPTX — zéro PNG, tout éditable dans Google Slides
"""

import os, re
from pptx import Presentation
from pptx.util import Cm, Pt, Emu, Mm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT_FILE = "/Users/julesgaveglio/vanzon-website-claude-code/public/docs/Guide-VBA-Lit-Peigne.pptx"
SVG_DIR = "/Volumes/T7/pieces_van_svg2"

# A4 Portrait
PAGE_W = Cm(21.01)
PAGE_H = Cm(29.69)

GREY = RGBColor(0x88, 0x88, 0x88)
LIGHT_GREY = RGBColor(0xDD, 0xDD, 0xDD)
FILL_GREY = RGBColor(0xDD, 0xDD, 0xDD)
BLACK = RGBColor(0x00, 0x00, 0x00)


# ─── Helpers ───

def add_text(slide, left, top, width, height, text, size=12, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    if size: p.font.size = Pt(size)
    if bold: p.font.bold = True
    if color: p.font.color.rgb = color
    p.alignment = align
    return txBox


def draw_polygon(slide, points_mm, offset_x, offset_y, scale):
    """Draw a native freeform polygon shape from mm coordinates.
    points_mm: list of (x_mm, y_mm) tuples
    Returns the shape (editable in Google Slides)
    """
    if not points_mm:
        return None

    # Calculate bounding box
    xs = [p[0] for p in points_mm]
    ys = [p[1] for p in points_mm]
    min_x, max_x = min(xs), max(xs)
    min_y, max_y = min(ys), max(ys)
    w = max_x - min_x
    h = max_y - min_y

    if w == 0 or h == 0:
        return None

    # Create freeform
    left = offset_x
    top = offset_y
    width = int(w * scale)
    height = int(h * scale)

    freeform = slide.shapes.build_freeform(left, top)

    # Move to first point
    first = points_mm[0]
    freeform.move_to(int((first[0] - min_x) * scale), int((first[1] - min_y) * scale))

    # Add line segments for remaining points
    segments = []
    for pt in points_mm[1:]:
        segments.append((int((pt[0] - min_x) * scale), int((pt[1] - min_y) * scale)))
    freeform.add_line_segments(segments, close=True)

    # Build
    shape = freeform.convert_to_shape(left, top)

    # Style
    shape.fill.solid()
    shape.fill.fore_color.rgb = FILL_GREY
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(1.2)

    return shape, width, height, min_x, min_y


def add_dim_line_h(slide, x1, y, x2, label, above=True):
    """Add a horizontal dimension line with arrows + label. All native shapes."""
    offset = Cm(0.8) if above else Cm(0.8)
    y_line = y - offset if above else y + offset
    y_text = y_line - Cm(0.5) if above else y_line + Cm(0.1)

    # Main line
    line = slide.shapes.add_connector(1, x1, y_line, x2, y_line)  # straight
    line.line.color.rgb = BLACK
    line.line.width = Pt(0.8)

    # End ticks
    tick_h = Cm(0.25)
    t1 = slide.shapes.add_connector(1, x1, y_line - tick_h, x1, y_line + tick_h)
    t1.line.color.rgb = BLACK; t1.line.width = Pt(0.8)
    t2 = slide.shapes.add_connector(1, x2, y_line - tick_h, x2, y_line + tick_h)
    t2.line.color.rgb = BLACK; t2.line.width = Pt(0.8)

    # Extension lines (dashed)
    ext1 = slide.shapes.add_connector(1, x1, y, x1, y_line)
    ext1.line.color.rgb = GREY; ext1.line.width = Pt(0.4)
    ext1.line.dash_style = 2  # dash
    ext2 = slide.shapes.add_connector(1, x2, y, x2, y_line)
    ext2.line.color.rgb = GREY; ext2.line.width = Pt(0.4)
    ext2.line.dash_style = 2

    # Label
    mid_x = (x1 + x2) // 2 - Cm(1)
    add_text(slide, mid_x, y_text, Cm(2), Cm(0.5), label, size=8, bold=True, align=PP_ALIGN.CENTER)


def add_dim_line_v(slide, x, y1, y2, label, left=True):
    """Add a vertical dimension line with label. All native shapes."""
    offset = Cm(0.8)
    x_line = x - offset if left else x + offset

    # Main line
    line = slide.shapes.add_connector(1, x_line, y1, x_line, y2)
    line.line.color.rgb = BLACK; line.line.width = Pt(0.8)

    # End ticks
    tick_w = Cm(0.25)
    t1 = slide.shapes.add_connector(1, x_line - tick_w, y1, x_line + tick_w, y1)
    t1.line.color.rgb = BLACK; t1.line.width = Pt(0.8)
    t2 = slide.shapes.add_connector(1, x_line - tick_w, y2, x_line + tick_w, y2)
    t2.line.color.rgb = BLACK; t2.line.width = Pt(0.8)

    # Extension lines
    ext1 = slide.shapes.add_connector(1, x, y1, x_line, y1)
    ext1.line.color.rgb = GREY; ext1.line.width = Pt(0.4); ext1.line.dash_style = 2
    ext2 = slide.shapes.add_connector(1, x, y2, x_line, y2)
    ext2.line.color.rgb = GREY; ext2.line.width = Pt(0.4); ext2.line.dash_style = 2

    # Label (rotated text not possible in pptx easily, so place horizontally)
    mid_y = (y1 + y2) // 2 - Cm(0.25)
    x_text = x_line - Cm(1.5) if left else x_line + Cm(0.2)
    add_text(slide, x_text, mid_y, Cm(1.3), Cm(0.5), label, size=8, bold=True, align=PP_ALIGN.CENTER)


def parse_svg_points(filename):
    """Parse polygon points from SVG file."""
    path = os.path.join(SVG_DIR, filename)
    if not os.path.exists(path):
        return [], 0, 0
    with open(path) as f:
        content = f.read()
    w = float(re.search(r'width="([\d.]+)mm"', content).group(1))
    h = float(re.search(r'height="([\d.]+)mm"', content).group(1))
    pts_m = re.search(r'points="([^"]+)"', content)
    coords = []
    if pts_m:
        for p in pts_m.group(1).strip().split(' '):
            xy = p.split(',')
            if len(xy) == 2:
                coords.append((float(xy[0]), float(xy[1])))
    return coords, w, h


def make_piece_page(slide_or_none, ref, name, dims_text, svg_file,
                    area_x, area_y, area_w, area_h, show_dims=True):
    """Draw a single piece as native shapes in a given area.
    If slide_or_none is None, creates a new slide.
    Returns the slide.
    """
    slide = slide_or_none if slide_or_none else prs.slides.add_slide(prs.slide_layouts[6])

    coords, svg_w, svg_h = parse_svg_points(svg_file)

    if not coords:
        # Simple rectangle fallback
        dims_parts = dims_text.split(" · ")[0].replace(",", ".").split(" × ")
        if len(dims_parts) >= 2:
            try:
                pw = float(dims_parts[0].replace(" cm", "").replace(" ", "")) * 10  # cm→mm
                ph = float(dims_parts[1].replace(" cm", "").replace(" ", "")) * 10
                coords = [(0,0), (pw,0), (pw,ph), (0,ph)]
                svg_w, svg_h = pw, ph
            except:
                return slide

    if not coords:
        return slide

    # Calculate scale to fit in area (with margin for dim lines)
    margin = Cm(1.5) if show_dims else Cm(0.3)
    avail_w = area_w - margin * 2
    avail_h = area_h - margin * 2

    xs = [p[0] for p in coords]
    ys = [p[1] for p in coords]
    shape_w = max(xs) - min(xs)
    shape_h = max(ys) - min(ys)

    if shape_w == 0 or shape_h == 0:
        return slide

    scale_x = avail_w / shape_w
    scale_y = avail_h / shape_h
    scale = min(scale_x, scale_y)

    # Position shape centered in area
    drawn_w = int(shape_w * scale)
    drawn_h = int(shape_h * scale)
    shape_x = area_x + margin + (avail_w - drawn_w) // 2
    shape_y = area_y + margin + (avail_h - drawn_h) // 2

    # Draw the polygon
    result = draw_polygon(slide, coords, shape_x, shape_y, scale)
    if not result:
        return slide

    _, actual_w, actual_h, min_x, min_y = result

    if show_dims:
        # Add dimension lines based on shape complexity
        n_pts = len(coords)

        if n_pts == 4:
            # Simple rectangle — width on top, height on left
            w_label = f"{shape_w:.0f}"
            h_label = f"{shape_h:.0f}"
            add_dim_line_h(slide, shape_x, shape_y, shape_x + drawn_w, w_label, above=True)
            add_dim_line_v(slide, shape_x, shape_y, shape_y + drawn_h, h_label, left=True)

        elif n_pts == 5:
            # Chamfered rectangle (A4 or A5)
            # Overall dims
            add_dim_line_h(slide, shape_x, shape_y + drawn_h, shape_x + drawn_w,
                          f"{shape_w:.0f}", above=False)
            add_dim_line_v(slide, shape_x, shape_y, shape_y + drawn_h,
                          f"{shape_h:.0f}", left=True)
            # Top edge (shorter)
            top_w = coords[1][0] - coords[0][0]
            top_drawn = int(top_w * scale)
            add_dim_line_h(slide, shape_x, shape_y, shape_x + top_drawn,
                          f"{top_w:.0f}", above=True)
            # Chamfer
            chamfer_x = shape_w - top_w
            chamfer_y = coords[2][1] - coords[1][1]
            add_text(slide, shape_x + top_drawn + Cm(0.1), shape_y + Cm(0.1),
                    Cm(2), Cm(0.4), f"chanfrein\n{chamfer_x:.0f}×{chamfer_y:.0f}", size=7, color=GREY)

        elif n_pts == 8:
            # Rectangle with notch (A1 planche sol, A2 fond)
            # Overall dims
            add_dim_line_h(slide, shape_x, shape_y, shape_x + drawn_w,
                          f"{shape_w:.0f}", above=True)
            add_dim_line_v(slide, shape_x, shape_y, shape_y + drawn_h,
                          f"{shape_h:.0f}", left=True)

            # Notch dimensions
            # Find the notch: where x or y changes unexpectedly
            # For planche sol: notch at (585,870)→(585,1705) carved from right
            if svg_file == "066_Planche sol.svg":
                notch_x = 585
                notch_y1 = 870
                notch_y2 = 1705
                notch_w = shape_w - notch_x
                notch_h = notch_y2 - notch_y1

                # Width before notch
                nx = shape_x + int(notch_x * scale)
                ny1 = shape_y + int(notch_y1 * scale)
                ny2 = shape_y + int(notch_y2 * scale)

                add_dim_line_h(slide, shape_x, shape_y + drawn_h, nx,
                              f"{notch_x:.0f}", above=False)
                add_dim_line_v(slide, shape_x + drawn_w, ny1, ny2,
                              f"{notch_h:.0f}", left=False)
                add_dim_line_h(slide, nx, ny1, shape_x + drawn_w,
                              f"{notch_w:.0f}", above=True)
                add_text(slide, nx + Cm(0.2), ny1 + Cm(0.2), Cm(2), Cm(0.4),
                        "encoche", size=7, color=GREY)

            elif svg_file == "067_Planche fond longueur.svg":
                # Notch: middle section narrows to 44mm wide
                notch_inner = 44
                notch_y1 = 870
                notch_y2 = 1705

                ni = shape_x + int(notch_inner * scale)
                ny1 = shape_y + int(notch_y1 * scale)
                ny2 = shape_y + int(notch_y2 * scale)

                add_dim_line_h(slide, shape_x, ny1 + (ny2-ny1)//2, ni,
                              f"{notch_inner:.0f}", above=True)
                add_dim_line_v(slide, shape_x + drawn_w, ny1, ny2,
                              f"{notch_y2-notch_y1:.0f}", left=False)

        else:
            # Complex shape (peigne, pied) — just bounding box
            add_dim_line_h(slide, shape_x, shape_y, shape_x + drawn_w,
                          f"{shape_w:.0f}", above=True)
            add_dim_line_v(slide, shape_x, shape_y, shape_y + drawn_h,
                          f"{shape_h:.0f}", left=True)
            if "eigne" in svg_file or "fixee" in svg_file:
                add_text(slide, shape_x + drawn_w//4, shape_y + drawn_h//2,
                        Cm(4), Cm(0.4), "Découpe depuis gabarit SVG", size=8, color=GREY)

    return slide


# ═══════════════════════════════════════
#  BUILD PRESENTATION
# ═══════════════════════════════════════
prs = Presentation()
prs.slide_width = PAGE_W
prs.slide_height = PAGE_H


def make_cover(title, dims, matiere, nb_pan, nb_car):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, Cm(0), Cm(7.6), Cm(21), Cm(2), title, size=28, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Cm(0), Cm(9.7), Cm(21), Cm(1.3), "Guide de montage", size=16, align=PP_ALIGN.CENTER)
    info = [("Encombrement fini", dims), ("Matière", matiere),
            ("Panneaux à découper", str(nb_pan)), ("Carrelets 15 × 15 mm", str(nb_car))]
    for i, (l, v) in enumerate(info):
        y = Cm(12.2 + i * 1.3)
        add_text(slide, Cm(4.1), y, Cm(6.6), Cm(1), l, size=11)
        add_text(slide, Cm(10.7), y, Cm(6.6), Cm(1), v, size=11)


def make_debit(title, planches, carrelets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, Cm(1.5), Cm(1.3), Cm(17.8), Cm(1.5), title, size=20, bold=True)
    y = Cm(3.3)
    if planches:
        rows = len(planches) + 1
        ts = slide.shapes.add_table(rows, 4, Cm(1.5), y, Cm(17.9), Cm(0.6 * rows))
        t = ts.table
        t.columns[0].width = Cm(2); t.columns[1].width = Cm(7.5)
        t.columns[2].width = Cm(6); t.columns[3].width = Cm(2.4)
        for j, h in enumerate(["Rép.", "Panneau (CP 15 mm)", "Dimensions", "Qté"]):
            t.cell(0, j).text = h
        for i, p in enumerate(planches):
            t.cell(i+1, 0).text = p[0]; t.cell(i+1, 1).text = p[1]
            t.cell(i+1, 2).text = p[2].split(" · ")[0]; t.cell(i+1, 3).text = "1"
        y += Cm(0.6 * rows + 1)
    if carrelets:
        rows2 = len(carrelets) + 1
        max_r = min(rows2, int((Cm(28) - y) / Cm(0.5)))
        ts2 = slide.shapes.add_table(max_r, 4, Cm(1.5), y, Cm(17.9), Cm(0.5 * max_r))
        t2 = ts2.table
        t2.columns[0].width = Cm(2); t2.columns[1].width = Cm(7.5)
        t2.columns[2].width = Cm(6); t2.columns[3].width = Cm(2.4)
        for j, h in enumerate(["Rép.", "Carrelet (15 × 15 mm)", "Longueur", "Qté"]):
            t2.cell(0, j).text = h
        for i, c in enumerate(carrelets[:max_r-1]):
            t2.cell(i+1, 0).text = c[0]; t2.cell(i+1, 1).text = c[1]
            t2.cell(i+1, 2).text = c[2]; t2.cell(i+1, 3).text = str(c[3])
        remaining = carrelets[max_r-1:]
        if remaining:
            make_debit(f"{title} (suite)", None, remaining)


def make_pieces_native(title, pieces):
    """One piece per half-page with native polygon + dimension lines."""
    items_per_slide = 2  # 2 pieces per page (top + bottom)

    for page in range(0, len(pieces), items_per_slide):
        batch = pieces[page:page + items_per_slide]
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        pn = page // items_per_slide + 1
        total = (len(pieces) + items_per_slide - 1) // items_per_slide
        sfx = f" ({pn}/{total})" if total > 1 else ""
        add_text(slide, Cm(1.5), Cm(0.5), Cm(17.8), Cm(1), f"Pièces à découper — {title}{sfx}", size=16, bold=True)

        for i, piece in enumerate(batch):
            ref, name, dims, svg = piece[0], piece[1], piece[2], piece[3]

            # Top half or bottom half
            y_start = Cm(2) + i * Cm(13.5)

            # Ref + name + dims text
            add_text(slide, Cm(1.5), y_start, Cm(1.5), Cm(0.6), ref, size=12, bold=True)
            add_text(slide, Cm(3), y_start, Cm(8), Cm(0.6), name, size=11)
            add_text(slide, Cm(1.5), y_start + Cm(0.6), Cm(15), Cm(0.5), dims, size=9, color=GREY)

            # Draw the piece with dimensions
            make_piece_page(slide, ref, name, dims, svg,
                          area_x=Cm(2.5), area_y=y_start + Cm(1.2),
                          area_w=Cm(16), area_h=Cm(12),
                          show_dims=True)


def make_steps(title, steps):
    steps_per_slide = 4
    for page in range(0, len(steps), steps_per_slide):
        batch = steps[page:page + steps_per_slide]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pn = page // steps_per_slide + 1
        total = (len(steps) + steps_per_slide - 1) // steps_per_slide
        sfx = f" ({pn}/{total})" if total > 1 else ""
        add_text(slide, Cm(1.5), Cm(1.3), Cm(17.8), Cm(1.5), f"Étapes de montage — {title}{sfx}", size=20, bold=True)
        for i, (num, st, instr) in enumerate(batch):
            y = Cm(3.6 + i * 3.3)
            add_text(slide, Cm(1.5), y, Cm(1.5), Cm(1.5), str(num), size=18, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, Cm(3.3), y, Cm(15.2), Cm(1), st, size=13, bold=True)
            add_text(slide, Cm(3.3), y + Cm(1), Cm(16), Cm(2), instr, size=10)


# ─── Piece data ───
BASE_PLANCHES = [
    ("A1", "Planche sol", "720,4 × 1870 mm · ép. 15 mm", "066_Planche sol.svg"),
    ("A2", "Planche fond longueur", "365 × 1870 mm · ép. 15 mm", "067_Planche fond longueur.svg"),
    ("A3", "Planche face longueur", "1870 × 365 mm · ép. 15 mm", "070_Planche face longueur.svg"),
    ("A4", "Planche droite largeur", "775 × 380 mm · ép. 15 mm", "068_Planche droite largeur.svg"),
    ("A5", "Planche gauche largeur", "380 × 775 mm · ép. 15 mm", "069_Planche gauche largeur.svg"),
]

BASE_CARRELETS = [
    ("B1", "Longueur haut gauche", "187 cm", 1),
    ("B2", "Longueur face bas", "187 cm", 1),
    ("B3", "Longueur fond bas droit", "84,5 cm", 1),
    ("B4", "Longueur fond bas gauche", "16,5 cm", 1),
    ("B5", "Largeur milieu haut", "73 cm", 1),
    ("B6", "Largeur milieu haut 2", "73 cm", 1),
    ("B7", "Largeur bas droite", "71,5 cm", 1),
    ("B8", "Largeur 1", "71,5 cm", 1),
    ("B9", "Largeur milieu bas 2", "71,5 cm", 1),
    ("B10", "Droite largeur 2", "65,8 cm", 1),
    ("B11", "Vertical avant droite", "26,5 cm", 1),
    ("B12", "Vertical avant gauche", "26,5 cm", 1),
    ("B13", "Vertical fond droite", "32 cm", 1),
    ("B14", "Vertical intérieur fond", "32 cm", 1),
    ("B15", "Gauche vertical", "32 cm", 1),
    ("B16", "Vertical fond gauche petit", "32 cm", 1),
    ("B17", "Bas gauche largeur petit", "14,5 cm", 1),
    ("B18", "Petit gauche haut", "14,5 cm", 1),
    ("B19", "Long. gauche support tiroir", "101 cm", 1),
    ("B20", "Long. bas gauche support tiroir", "101 cm", 1),
    ("B21", "Long. droite support tiroir", "101 cm", 1),
    ("B22", "Long. bas droit support tiroir", "101 cm", 1),
    ("B23", "Vert. fond gauche support tiroir", "32 cm", 1),
    ("B24", "Vert. fond droit support tiroir", "33,5 cm", 1),
]

CLOISONS = [
    ("C1", "Support gauche glissière", "1010 × 350 mm · ép. 15 mm", "093_Planche support gauche glissi_re.svg"),
    ("C2", "Planche droite intérieur", "1010 × 350 mm · ép. 15 mm", "095_Planche droite int_rieur.svg"),
    ("C3", "Planche placard gauche", "745 × 350 mm · ép. 15 mm", "094_Planche placard gauche.svg"),
    ("C4", "Séparation boxio / rangement", "745 × 350 mm · ép. 15 mm", "096_Planche s_paration boxio _ rangement.svg"),
]

TIROIR_PLANCHES = [
    ("D1", "Planche bas tiroir", "985 × 410 mm · ép. 15 mm", "083_Planche bas cuisine.svg"),
    ("D2", "Planche fond tiroir", "440 × 310 mm · ép. 15 mm", "084_Planche fond cuisine.svg"),
    ("D3", "Planche face tiroir", "440 × 310 mm · ép. 15 mm", "087_Planche cuisine face.svg"),
    ("D4", "Planche droite tiroir", "985 × 310 mm · ép. 15 mm", "085_Planche droite cuisine.svg"),
    ("D5", "Planche gauche tiroir", "985 × 310 mm · ép. 15 mm", "086_Planche gauche cuisine.svg"),
]

TIROIR_CARRELETS = [
    ("D6", "Long. gauche tiroir", "98,5 cm", 1),
    ("D7", "Long. droit tiroir", "98,5 cm", 1),
    ("D8", "Fond gauche vertical", "28 cm", 1),
    ("D9", "Fond droit vertical", "28 cm", 1),
    ("D10", "Vertical avant gauche", "28 cm", 1),
    ("D11", "Vertical avant droit", "28 cm", 1),
    ("D12", "Fond largeur", "38 cm", 1),
    ("D13", "Largeur bas face", "38 cm", 1),
]

PEIGNE_PLANCHES = [
    ("E1", "Planche du dessus lit", "775 × 1870 mm · ép. 15 mm", "056_Planche du dessus lit.svg"),
    ("E2", "Lit peigne fixée", "740 × 1900 mm · ép. 15 mm", "057_Lit peigne fixee.svg"),
    ("E3", "Peigne mobile droite", "1032 × 845 mm · ép. 15 mm", "061_Peigne mobile partie droite.svg"),
    ("E4", "Peigne mobile gauche", "950 × 740 mm · ép. 15 mm", "062_Peigne mobile partie gauche.svg"),
    ("E5", "Pied lit droite", "380 × 950 mm · ép. 15 mm", "060_Pied lit peigne partie droite.svg"),
    ("E6", "Pied lit gauche", "380 × 950 mm · ép. 15 mm", "063_Pied lit peigne partie gauche.svg"),
    ("E7", "Triangle pied droit", "100 × 100 mm · ép. 15 mm", "059_Triangle pied droit.svg"),
    ("E8", "Triangle pied gauche", "100 × 100 mm · ép. 15 mm", "065_Triangle pied gauche.svg"),
]


# ═══════════════════════════════════════
#  GENERATE ALL SLIDES
# ═══════════════════════════════════════

# BLOC 1
make_cover("Bloc Lit Peigne\nBase extérieure", "187 × 77,5 × 38 cm", "CP peuplier 15 mm", 5, 24)
make_debit("Liste de débit — Base extérieure", BASE_PLANCHES, BASE_CARRELETS)
make_pieces_native("Base extérieure", BASE_PLANCHES)
make_steps("Base extérieure", [
    (1, "Poser la planche sol [A1]", "Place A1 à plat. Attention : encoche passage de roue — vérifie l'orientation."),
    (2, "Monter les côtés [A4] et [A5]", "Visse A4 et A5 verticalement sur les bords de A1. Pré-perce à 3 mm, 4 vis 4×40 mm par côté."),
    (3, "Fixer le fond [A2]", "Glisse A2 à l'arrière. Visse dans les chants des côtés — 4 vis 3,5×30 mm par côté."),
    (4, "Fermer avec la face [A3]", "Fixe A3 à l'avant. Vérifie l'équerrage (diagonales) avant de serrer."),
])
make_steps("Ossature carrelets", [
    (5, "Cadre bas [B2, B3, B4, B7-B9]", "Visse les carrelets horizontaux au fond. B2 (face, 187 cm), B3+B4 (fond), B7-B9 (largeurs). 2 vis 3,5×25 mm."),
    (6, "Montants verticaux [B11-B16]", "Fixe les 6 carrelets verticaux dans les angles. Visse par le bas + par le côté."),
    (7, "Cadre haut [B1, B5, B6, B10]", "Pose les carrelets du haut. B1 (187 cm), B5+B6 (73 cm), B10 (65,8 cm). 2 vis par jonction."),
    (8, "Petits carrelets + rails [B17-B24]", "B17-B18 (14,5 cm) ferment le cadre gauche. B19-B24 forment les rails du tiroir (parallèles, ~101 cm)."),
])

# BLOC 3
make_cover("Bloc Lit Peigne\nCloisons intérieures", "—", "CP peuplier 15 mm", 4, 0)
make_debit("Liste de débit — Cloisons", CLOISONS)
make_pieces_native("Cloisons", CLOISONS)
make_steps("Cloisons intérieures", [
    (1, "Cloisons longitudinales [C1] et [C2]", "Visse C1 et C2 verticalement sur les carrelets. 4 vis 3,5×30 mm par planche."),
    (2, "Cloisons transversales [C3] et [C4]", "Fixe C3 et C4 perpendiculairement. Visse dans carrelets + cloisons longitudinales."),
])

# BLOC 4
make_cover("Bloc Lit Peigne\nTiroir coulissant", "98,5 × 41 × 31 cm", "CP peuplier 15 mm", 5, 8)
make_debit("Liste de débit — Tiroir", TIROIR_PLANCHES, TIROIR_CARRELETS)
make_pieces_native("Tiroir — planches", TIROIR_PLANCHES)
make_steps("Tiroir coulissant", [
    (1, "Ossature tiroir [D6-D13]", "Assemble les 8 carrelets en cadre : 2 longueurs + 2 largeurs + 4 verticaux."),
    (2, "Planches du caisson [D1-D5]", "Visse D1 (bas), puis D2-D5 (4 côtés). 3 vis par planche."),
    (3, "Glissières + insertion", "Fixe les glissières sur les rails B19-B22. Glisse le tiroir. Il doit coulisser librement."),
])

# BLOC 5
make_cover("Bloc Lit Peigne\nPeigne — surface de couchage", "190 × 152 cm (déplié)", "CP peuplier 15 mm", 8, 0)
make_debit("Liste de débit — Peigne", PEIGNE_PLANCHES)
make_pieces_native("Peigne", PEIGNE_PLANCHES)
make_steps("Peigne (lit)", [
    (1, "Pieds + triangles [E5-E8]", "Colle et visse les triangles E7-E8 sur les pieds E5-E6. 2 vis 3,5×25 mm + colle."),
    (2, "Pieds sur peignes mobiles [E3, E4]", "Visse chaque pied perpendiculairement au peigne mobile. 6 vis 3,5×30 mm par pied."),
    (3, "Planche du dessus [E1]", "Place E1 sur le caisson, sur les carrelets haut. 6 vis 3,5×30 mm."),
    (4, "Peigne fixe [E2]", "Pose E2 par-dessus E1, dents vers le haut. 4 vis 3,5×30 mm."),
    (5, "Peignes mobiles", "E3 et E4 se replient sur le fixe. Dépliés, les dents s'emboîtent. Teste le dépliage."),
])


# ═══════════════════════════════════════
prs.save(OUT_FILE)
print(f"✅ PPTX v3 : {OUT_FILE}")
print(f"   Format : Portrait A4")
print(f"   Slides : {len(prs.slides)}")
print(f"   100% formes natives — zéro PNG")
print(f"   Taille : {os.path.getsize(OUT_FILE) / 1024:.0f} KB")
