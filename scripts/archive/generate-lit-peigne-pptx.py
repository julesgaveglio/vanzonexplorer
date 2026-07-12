#!/usr/bin/env python3
"""
Génère le guide de montage PPTX du Bloc Lit Peigne — style IKEA / Guide VBA Construction
Usage: DYLD_LIBRARY_PATH=/opt/homebrew/lib python3 scripts/generate-lit-peigne-pptx.py
"""

import os
import cairosvg
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Config ───
SVG_DIR = "/Volumes/T7/pieces_van_svg2"
OUT_DIR = "/Users/julesgaveglio/vanzon-website-claude-code/public/docs"
PNG_CACHE = os.path.join(OUT_DIR, "png-cache-lit")
OUT_FILE = os.path.join(OUT_DIR, "Guide-VBA-Construction-Lit-Peigne.pptx")

os.makedirs(PNG_CACHE, exist_ok=True)

# ─── Couleurs ───
BLUE = RGBColor(0x2B, 0x5E, 0xA7)
DARK = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xB9, 0x94, 0x5F)
WARN_BG = RGBColor(0xFF, 0xF8, 0xE1)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)

# ─── Slide size (16:9) ───
SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)

# ─── SVG → PNG helper ───
def svg_to_png(svg_filename, max_width=300):
    """Convert SVG to PNG, return path. Cached."""
    png_path = os.path.join(PNG_CACHE, svg_filename.replace(".svg", ".png"))
    if not os.path.exists(png_path):
        svg_path = os.path.join(SVG_DIR, svg_filename)
        if os.path.exists(svg_path):
            cairosvg.svg2png(url=svg_path, write_to=png_path, output_width=max_width)
    return png_path if os.path.exists(png_path) else None


# ─── Text helpers ───
def add_text(slide, left, top, width, height, text, size=10, bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name="Inter"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_rich_text(slide, left, top, width, height):
    """Return a text frame for multi-run paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox.text_frame


def add_rounded_rect(slide, left, top, width, height, fill_color=LIGHT_GREY, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_circle_number(slide, left, top, number, size=Cm(1.2), bg_color=BLUE, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = "Inter"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)


# ─── Piece data ───
# Bloc 1: Base extérieure (caisson)
BASE_PLANCHES = [
    ("A1", "Planche sol", "187 × 74,5 cm", "066_Planche sol.svg", 1, "Forme avec encoche passage de roue"),
    ("A2", "Planche fond (longueur)", "187 × 36,5 cm", "067_Planche fond longueur.svg", 1, None),
    ("A3", "Planche face (longueur)", "187 × 36,5 cm", "070_Planche face longueur.svg", 1, None),
    ("A4", "Planche droite (largeur)", "77,5 × 38 cm", "068_Planche droite largeur.svg", 1, None),
    ("A5", "Planche gauche (largeur)", "77,5 × 38 cm", "069_Planche gauche largeur.svg", 1, None),
]

# Bloc 2: Ossature carrelets
BASE_CARRELETS = [
    ("B1", "Carrelet longueur haut gauche", "187 cm", "102_Tasseau longueur haut gauche.svg", 1),
    ("B2", "Carrelet longueur face bas", "187 cm", "111_Tasseau longueur face bas.svg", 1),
    ("B3", "Carrelet longueur fond bas droit", "84,5 cm", "112_Tasseau longueur fond bas droit.svg", 1),
    ("B4", "Carrelet longueur fond bas gauche", "16,5 cm", "113_Tasseau longueur fond bas gauche.svg", 1),
    ("B5", "Carrelet largeur milieu haut", "73 cm", "103_Tasseau largeur milieu haut.svg", 1),
    ("B6", "Carrelet largeur milieu haut 2", "73 cm", "120_Tasseau largeur milieu haut 2.svg", 1),
    ("B7", "Carrelet largeur bas droite", "71,5 cm", "104_Tasseau largeur bas droite.svg", 1),
    ("B8", "Carrelet largeur 1", "71,5 cm", "110_Tasseau largeur 1.svg", 1),
    ("B9", "Carrelet largeur milieu bas 2", "71,5 cm", "119_Tasseau largeur milieu bas 2.svg", 1),
    ("B10", "Carrelet droite largeur 2", "65,8 cm", "118_Tasseau droite largeur 2.svg", 1),
    ("B11", "Carrelet vertical avant droite", "26,5 cm", "105_Tasseau vertical avant droite.svg", 1),
    ("B12", "Carrelet vertical avant gauche", "26,5 cm", "109_Tasseau vertical avant gauche 1.svg", 1),
    ("B13", "Carrelet vertical fond droite", "32 cm", "101_Tasseau vertical fond droite.svg", 1),
    ("B14", "Carrelet vertical intérieur fond", "32 cm", "108_Tasseau vertical int_rieur fond.svg", 1),
    ("B15", "Carrelet gauche vertical", "32 cm", "115_Tasseau gauche vertical.svg", 1),
    ("B16", "Carrelet vertical fond gauche petit", "32 cm", "117_tasseau vertical fond gauche petit.svg", 1),
    ("B17", "Carrelet bas gauche largeur petit", "14,5 cm", "114_Tasseau bas gauche largeur petit.svg", 1),
    ("B18", "Carrelet petit gauche haut", "14,5 cm", "116_Tasseau petit gauche haut.svg", 1),
]

# Bloc 2b: Carrelets support tiroir
TIROIR_CARRELETS = [
    ("B19", "Carrelet longueur gauche support tiroir", "101 cm", "097_Tasseau longueur gauche support tiroir.svg", 1),
    ("B20", "Carrelet longueur bas gauche support tiroir", "101 cm", "099_Tasseau longueur bas gauche support tiroir.svg", 1),
    ("B21", "Carrelet longueur droite support tiroir", "101 cm", "107_Tasseau longueur droite support tiroir.svg", 1),
    ("B22", "Carrelet longueur bas droit support tiroir", "101 cm", "106_Tasseau longueur bas droit support tiroir.svg", 1),
    ("B23", "Carrelet vertical fond gauche support tiroir", "32 cm", "098_Tasseau vertical fond gauche support tiroir.svg", 1),
    ("B24", "Carrelet vertical fond droit support tiroir", "33,5 cm", "100_Tasseau vertical fond droit support tiroir.svg", 1),
]

# Bloc 3: Cloisons intérieures
CLOISONS = [
    ("C1", "Planche support gauche glissière", "101 × 35 cm", "093_Planche support gauche glissi_re.svg", 1, None),
    ("C2", "Planche droite intérieur", "101 × 35 cm", "095_Planche droite int_rieur.svg", 1, None),
    ("C3", "Planche placard gauche", "74,5 × 35 cm", "094_Planche placard gauche.svg", 1, None),
    ("C4", "Planche séparation boxio / rangement", "74,5 × 35 cm", "096_Planche s_paration boxio _ rangement.svg", 1, None),
]

# Bloc 4: Tiroir coulissant (caisson)
TIROIR_PLANCHES = [
    ("D1", "Planche bas tiroir", "98,5 × 41 cm", "083_Planche bas cuisine.svg", 1, None),
    ("D2", "Planche fond tiroir", "44 × 31 cm", "084_Planche fond cuisine.svg", 1, None),
    ("D3", "Planche face tiroir", "44 × 31 cm", "087_Planche cuisine face.svg", 1, None),
    ("D4", "Planche droite tiroir", "98,5 × 31 cm", "085_Planche droite cuisine.svg", 1, None),
    ("D5", "Planche gauche tiroir", "98,5 × 31 cm", "086_Planche gauche cuisine.svg", 1, None),
]

TIROIR_CARRELETS_INT = [
    ("D6", "Carrelet longueur gauche tiroir", "98,5 cm", "076_Tasseau longueur gauche tiroir.svg", 1),
    ("D7", "Carrelet longueur droit tiroir", "98,5 cm", "078_Tasseau longueur droit tiroir.svg", 1),
    ("D8", "Carrelet fond gauche vertical", "28 cm", "071_Tasseau fond gauche vertical.svg", 1),
    ("D9", "Carrelet fond droit vertical", "28 cm", "073_Tasseau fond droit vertical.svg", 1),
    ("D10", "Carrelet vertical avant gauche", "28 cm", "074_Tasseau vertical avant gauche.svg", 1),
    ("D11", "Carrelet vertical avant droit", "28 cm", "077_Tasseau vertical avant droit.svg", 1),
    ("D12", "Carrelet fond largeur", "38 cm", "072_Tasseau fond largeur.svg", 1),
    ("D13", "Carrelet largeur bas face", "38 cm", "075_Tasseau largeur bas face.svg", 1),
]

# Bloc 5: Peigne (lit)
PEIGNE = [
    ("E1", "Planche du dessus lit", "187 × 77,5 cm", "056_Planche du dessus lit.svg", 1, None),
    ("E2", "Lit peigne fixée", "190 × 74 cm", "057_Lit peigne fixee.svg", 1, "Se pose sur le dessus"),
    ("E3", "Peigne mobile partie droite", "95 × 74 cm", "061_Peigne mobile partie droite.svg", 1, "Forme avec dents de peigne"),
    ("E4", "Peigne mobile partie gauche", "95 × 74 cm", "062_Peigne mobile partie gauche.svg", 1, "Forme avec dents de peigne"),
    ("E5", "Pied lit peigne partie droite", "95 × 38 cm", "060_Pied lit peigne partie droite.svg", 1, None),
    ("E6", "Pied lit peigne partie gauche", "95 × 38 cm", "063_Pied lit peigne partie gauche.svg", 1, None),
    ("E7", "Triangle pied droit", "10 × 10 cm", "059_Triangle pied droit.svg", 1, "Renfort"),
    ("E8", "Triangle pied gauche", "10 × 10 cm", "065_Triangle pied gauche.svg", 1, "Renfort"),
]

# Accessoires (achetés, non construits)
ACCESSOIRES = [
    ("—", "Glissière partie fixe droite", "100 cm", "091"),
    ("—", "Glissière partie fixe gauche", "100 cm", "092"),
    ("—", "Glissière coulissante droite", "103,3 cm", "082"),
    ("—", "Glissière coulissante gauche", "103,3 cm", "081"),
    ("—", "Support table extérieur (×2 pièces)", "40 × 5 cm", "079-080"),
    ("—", "Table extérieur", "90 × 40 cm", "090"),
    ("—", "Boxio gaz", "—", "088"),
    ("—", "Boxio robinet", "—", "089"),
]


# ═══════════════════════════════════════
#  BUILD PRESENTATION
# ═══════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ─── HELPER: Slide cover pour un bloc ───
def make_cover_slide(title, subtitle, dimensions, nb_planches, nb_carrelets, bloc_num, total_blocs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Bloc number badge
    add_rounded_rect(slide, Cm(1.5), Cm(1), Cm(5), Cm(1.2), fill_color=BLUE)
    add_text(slide, Cm(1.5), Cm(1), Cm(5), Cm(1.2),
             f"BLOC {bloc_num}/{total_blocs}", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Title
    add_text(slide, Cm(1.5), Cm(3), Cm(20), Cm(2.5),
             title, size=32, bold=True, color=BLUE, font_name="Inter")

    # Subtitle
    add_text(slide, Cm(1.5), Cm(5.5), Cm(20), Cm(1.5),
             subtitle, size=14, color=GREY)

    # Info cards
    cards = [
        ("Encombrement", dimensions),
        ("Matière", "CP peuplier 15 mm"),
        ("Panneaux à découper", str(nb_planches)),
        ("Carrelets 15 × 15 mm", str(nb_carrelets)),
    ]

    card_w = Cm(7)
    card_h = Cm(3.5)
    start_x = Cm(1.5)
    y = Cm(9)
    gap = Cm(0.5)

    for i, (label, value) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rounded_rect(slide, x, y, card_w, card_h, fill_color=LIGHT_GREY)
        add_text(slide, x, y + Cm(0.4), card_w, Cm(1),
                 label, size=9, color=GREY, align=PP_ALIGN.CENTER)
        add_text(slide, x, y + Cm(1.4), card_w, Cm(1.5),
                 value, size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # Footer
    add_text(slide, Cm(1.5), Cm(17.5), Cm(30), Cm(1),
             "Guide VBA Construction — Bloc Lit Peigne — vanzonexplorer.com",
             size=8, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)

    return slide


# ─── HELPER: Slide liste de débit ───
def make_debit_slide(title, planches, carrelets=None, page_note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE

    add_text(slide, Cm(1.5), Cm(0.8), Cm(20), Cm(1.5),
             f"Liste de débit — {title}", size=22, bold=True, color=BLUE)

    # Table planches
    if planches:
        cols = 4
        rows = len(planches) + 1
        tbl_w = Cm(30)
        tbl_h = Cm(0.8) * rows
        y_start = Cm(2.5)

        table_shape = slide.shapes.add_table(rows, cols, Cm(1.5), y_start, tbl_w, tbl_h)
        table = table_shape.table

        # Column widths
        table.columns[0].width = Cm(3)
        table.columns[1].width = Cm(14)
        table.columns[2].width = Cm(8)
        table.columns[3].width = Cm(5)

        # Header
        headers = ["Rép.", "Panneau (CP 15 mm)", "Dimensions", "Qté"]
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
                paragraph.font.name = "Inter"
                paragraph.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE

        # Data
        for i, piece in enumerate(planches):
            ref, name, dims = piece[0], piece[1], piece[2]
            qty = piece[4] if len(piece) > 4 else 1
            row_data = [ref, name, dims, str(qty)]
            for j, val in enumerate(row_data):
                cell = table.cell(i + 1, j)
                cell.text = val
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(9)
                    paragraph.font.color.rgb = DARK
                    paragraph.font.name = "Inter"
                    paragraph.alignment = PP_ALIGN.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_GREY

    # Table carrelets
    if carrelets:
        tbl_w = Cm(30)
        if planches:
            y_carr = y_start + tbl_h + Cm(0.8)
        else:
            y_carr = Cm(2.5)

        cols2 = 4
        rows2 = len(carrelets) + 1
        tbl_h2 = Cm(0.7) * rows2

        # Check if fits on slide, if not just show what fits
        max_rows = min(rows2, int((Cm(18) - y_carr) / Cm(0.7)))
        if max_rows < rows2:
            rows2 = max_rows

        table_shape2 = slide.shapes.add_table(rows2, cols2, Cm(1.5), y_carr, tbl_w, tbl_h2)
        table2 = table_shape2.table

        table2.columns[0].width = Cm(3)
        table2.columns[1].width = Cm(14)
        table2.columns[2].width = Cm(8)
        table2.columns[3].width = Cm(5)

        headers2 = ["Rép.", "Carrelet (15 × 15 mm)", "Longueur", "Qté"]
        for j, h in enumerate(headers2):
            cell = table2.cell(0, j)
            cell.text = h
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
                paragraph.font.name = "Inter"
                paragraph.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = GOLD

        for i, piece in enumerate(carrelets[:rows2-1]):
            ref, name, length = piece[0], piece[1], piece[2]
            qty = piece[4] if len(piece) > 4 else 1
            row_data = [ref, name, length, str(qty)]
            for j, val in enumerate(row_data):
                cell = table2.cell(i + 1, j)
                cell.text = val
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(8)
                    paragraph.font.color.rgb = DARK
                    paragraph.font.name = "Inter"
                    paragraph.alignment = PP_ALIGN.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_GREY

    if page_note:
        add_text(slide, Cm(1.5), Cm(17.5), Cm(30), Cm(1),
                 page_note, size=8, color=GREY, align=PP_ALIGN.LEFT)

    return slide


# ─── HELPER: Slide pièces à découper ───
def make_pieces_slide(title, pieces, is_carrelet=False):
    """Creates piece cards with SVG thumbnails. Max ~6 per slide."""
    slides_created = []
    items_per_slide = 6

    for page in range(0, len(pieces), items_per_slide):
        batch = pieces[page:page + items_per_slide]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE

        page_num = page // items_per_slide + 1
        total_pages = (len(pieces) + items_per_slide - 1) // items_per_slide
        suffix = f" ({page_num}/{total_pages})" if total_pages > 1 else ""

        add_text(slide, Cm(1.5), Cm(0.5), Cm(25), Cm(1.2),
                 f"Pièces à découper — {title}{suffix}", size=18, bold=True, color=BLUE)

        # Placeholder for exploded view screenshot
        if page == 0:
            add_rounded_rect(slide, Cm(22), Cm(0.3), Cm(10.5), Cm(6), fill_color=LIGHT_GREY)
            add_text(slide, Cm(22), Cm(2.5), Cm(10.5), Cm(1.5),
                     "Coller ici le screenshot\nde la vue éclatée", size=9, color=GREY, align=PP_ALIGN.CENTER)

        # Piece cards (2 columns × 3 rows)
        card_w = Cm(15)
        card_h = Cm(2.5)
        col1_x = Cm(1.5)
        col2_x = Cm(17)
        start_y = Cm(2.2) if page > 0 else Cm(2.2)
        gap_y = Cm(0.3)

        for i, piece in enumerate(batch):
            col = i % 2
            row = i // 2
            x = col1_x if col == 0 else col2_x
            y = start_y + row * (card_h + gap_y)

            # Card background
            add_rounded_rect(slide, x, y, card_w, card_h, fill_color=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD))

            ref = piece[0]
            name = piece[1]
            dims = piece[2]
            svg_file = piece[3]
            note = piece[5] if len(piece) > 5 and piece[5] else None

            # Ref badge
            badge_color = GOLD if is_carrelet else BLUE
            add_rounded_rect(slide, x + Cm(0.3), y + Cm(0.5), Cm(1.8), Cm(1.5), fill_color=badge_color)
            add_text(slide, x + Cm(0.3), y + Cm(0.8), Cm(1.8), Cm(1),
                     ref, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

            # Name + dims
            add_text(slide, x + Cm(2.5), y + Cm(0.3), Cm(7), Cm(1),
                     name, size=10, bold=True, color=DARK)
            dim_label = dims
            if note:
                dim_label += f"  ·  {note}"
            add_text(slide, x + Cm(2.5), y + Cm(1.2), Cm(7), Cm(1),
                     dim_label, size=8, color=GREY)

            # SVG thumbnail (skip if SKIP_IMAGES env is set)
            if not os.environ.get('SKIP_IMAGES'):
                png_path = svg_to_png(svg_file, max_width=200)
                if png_path:
                    try:
                        thumb_w = Cm(4.5)
                        thumb_h = Cm(2)
                        slide.shapes.add_picture(png_path, x + card_w - thumb_w - Cm(0.3), y + Cm(0.25), thumb_w, thumb_h)
                    except Exception:
                        pass

        slides_created.append(slide)

    return slides_created


# ─── HELPER: Slide étapes de montage ───
def make_steps_slide(title, steps):
    """Steps: list of (step_num, title, instructions)"""
    slides_created = []
    steps_per_slide = 4

    for page in range(0, len(steps), steps_per_slide):
        batch = steps[page:page + steps_per_slide]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE

        page_num = page // steps_per_slide + 1
        total_pages = (len(steps) + steps_per_slide - 1) // steps_per_slide
        suffix = f" ({page_num}/{total_pages})" if total_pages > 1 else ""

        add_text(slide, Cm(1.5), Cm(0.5), Cm(25), Cm(1.5),
                 f"Étapes de montage — {title}{suffix}", size=18, bold=True, color=BLUE)

        step_h = Cm(3.5)
        start_y = Cm(2.2)

        for i, (num, step_title, instructions) in enumerate(batch):
            y = start_y + i * (step_h + Cm(0.3))

            # Step number circle
            add_circle_number(slide, Cm(1.5), y + Cm(0.2), num)

            # Step title
            add_text(slide, Cm(3.2), y, Cm(28), Cm(1),
                     step_title, size=13, bold=True, color=DARK)

            # Instructions
            add_text(slide, Cm(3.2), y + Cm(1.2), Cm(28), Cm(2.2),
                     instructions, size=10, color=GREY)

            # Separator line
            if i < len(batch) - 1:
                line_y = y + step_h + Cm(0.1)
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3.2), line_y, Cm(28), Pt(0.5))
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
                line.line.fill.background()

        slides_created.append(slide)

    return slides_created


# ═══════════════════════════════════════
#  SLIDE 1 — COUVERTURE GÉNÉRALE
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE

add_text(slide, Cm(1.5), Cm(2), Cm(30), Cm(2),
         "VANZON", size=40, bold=True, color=BLUE, align=PP_ALIGN.CENTER, font_name="Inter")
add_text(slide, Cm(1.5), Cm(4.5), Cm(30), Cm(1),
         "E X P L O R E R", size=12, color=GREY, align=PP_ALIGN.CENTER)

# Separator line
sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(12), Cm(6.5), Cm(10), Pt(1.5))
sep.fill.solid(); sep.fill.fore_color.rgb = BLUE; sep.line.fill.background()

add_text(slide, Cm(1.5), Cm(8), Cm(30), Cm(2),
         "Guide de montage", size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(slide, Cm(1.5), Cm(10.5), Cm(30), Cm(1.5),
         "Bloc Lit Peigne", size=18, color=GREY, align=PP_ALIGN.CENTER)

# Info cards
info_cards = [
    ("📦", "5 blocs"),
    ("📐", "~30 planches"),
    ("🔩", "~24 carrelets"),
    ("⏱️", "~6 heures"),
]
for i, (icon, label) in enumerate(info_cards):
    x = Cm(4) + i * Cm(7)
    add_text(slide, x, Cm(13), Cm(5), Cm(1.5),
             icon, size=24, align=PP_ALIGN.CENTER)
    add_text(slide, x, Cm(14.8), Cm(5), Cm(1),
             label, size=9, color=GREY, align=PP_ALIGN.CENTER)

add_text(slide, Cm(1.5), Cm(17.5), Cm(30), Cm(1),
         "v1.0 — Juin 2026", size=8, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
#  BLOC 1 — BASE EXTÉRIEURE
# ═══════════════════════════════════════
make_cover_slide(
    "Base extérieure",
    "Caisson principal — 5 planches CP peuplier formant la structure du lit",
    "187 × 77,5 × 38 cm",
    nb_planches=5, nb_carrelets=0, bloc_num=1, total_blocs=5
)

make_debit_slide("Base extérieure", BASE_PLANCHES)
make_pieces_slide("Base extérieure", BASE_PLANCHES)
make_steps_slide("Base extérieure", [
    (1, "Poser la planche sol [A1]",
     "Place la planche sol à plat sur ton plan de travail. Attention : cette planche a une encoche pour le passage de roue — vérifie l'orientation avant de continuer."),
    (2, "Monter les côtés [A4] et [A5]",
     "Visse les planches droite et gauche verticalement sur les bords latéraux de la planche sol. Pré-perce à 3 mm, puis 4 vis CP 4×40 mm par côté depuis le dessous."),
    (3, "Fixer le fond [A2]",
     "Glisse la planche fond longueur à l'arrière entre les deux côtés. Visse dans les chants des côtés — 4 vis 3,5×30 mm par côté + 3 vis dans le chant arrière de la planche sol."),
    (4, "Fermer avec la face [A3]",
     "Fixe la planche face longueur à l'avant, symétrique au fond. Mêmes vis que l'étape 3. Vérifie l'équerrage (diagonales) avant de serrer définitivement."),
])


# ═══════════════════════════════════════
#  BLOC 2 — OSSATURE CARRELETS
# ═══════════════════════════════════════
make_cover_slide(
    "Ossature carrelets",
    "Structure interne — 24 carrelets 15×15 mm qui rigidifient le caisson",
    "187 × 74,5 × 35 cm",
    nb_planches=0, nb_carrelets=24, bloc_num=2, total_blocs=5
)

make_debit_slide("Ossature (carrelets structure)", None, BASE_CARRELETS)
make_debit_slide("Ossature (carrelets support tiroir)", None, TIROIR_CARRELETS,
                 page_note="Ces 6 carrelets forment les rails sur lesquels le tiroir coulissant glisse.")

make_pieces_slide("Ossature carrelets",
    [(c[0], c[1], c[2], c[3], c[4], None) for c in BASE_CARRELETS], is_carrelet=True)
make_pieces_slide("Carrelets support tiroir",
    [(c[0], c[1], c[2], c[3], c[4], None) for c in TIROIR_CARRELETS], is_carrelet=True)

make_steps_slide("Ossature carrelets", [
    (1, "Cadre bas — carrelets de fond [B2, B3, B4, B8, B9]",
     "Visse les carrelets horizontaux au fond de la base, le long des bords internes. B2 (face, 187 cm), B3+B4 (fond), B8+B9 (largeurs). 2 vis 3,5×25 mm par carrelet."),
    (2, "Montants verticaux [B11-B16]",
     "Fixe les 6 carrelets verticaux dans les angles et jonctions internes. Ils portent les carrelets du haut. Visse par le bas dans le cadre + par le côté dans les planches."),
    (3, "Cadre haut [B1, B5, B6, B10]",
     "Pose les carrelets horizontaux du haut sur les montants. B1 (187 cm, longueur gauche), B5+B6 (largeurs milieu), B10 (largeur droite). 2 vis par jonction."),
    (4, "Petits carrelets de renfort [B17, B18]",
     "Fixe les 2 petits carrelets (14,5 cm) en bas à gauche (B17) et en haut à gauche (B18). Ils ferment le cadre dans la zone réduite."),
    (5, "Rails support tiroir [B19-B24]",
     "Visse les 4 carrelets horizontaux (B19-B22, ~101 cm) parallèles deux à deux — ils forment les rails du tiroir. Puis les 2 montants verticaux (B23-B24) à l'arrière. Vérifie le parallélisme : le tiroir doit coulisser sans forcer."),
])


# ═══════════════════════════════════════
#  BLOC 3 — CLOISONS INTÉRIEURES
# ═══════════════════════════════════════
make_cover_slide(
    "Cloisons intérieures",
    "4 planches verticales qui créent les compartiments (rangement, boxio, tiroir)",
    "—",
    nb_planches=4, nb_carrelets=0, bloc_num=3, total_blocs=5
)

make_debit_slide("Cloisons intérieures", CLOISONS)
make_pieces_slide("Cloisons intérieures", CLOISONS)
make_steps_slide("Cloisons intérieures", [
    (1, "Cloisons longitudinales [C1] et [C2]",
     "Visse les planches C1 (support glissière gauche) et C2 (droite intérieur) verticalement sur les carrelets. Elles séparent le tiroir des zones latérales. 4 vis 3,5×30 mm par planche."),
    (2, "Cloisons transversales [C3] et [C4]",
     "Fixe la planche placard gauche [C3] et la séparation boxio [C4] perpendiculairement. C4 délimite l'espace boxio du rangement. Visse dans les carrelets + dans les cloisons longitudinales."),
])


# ═══════════════════════════════════════
#  BLOC 4 — TIROIR COULISSANT
# ═══════════════════════════════════════
make_cover_slide(
    "Tiroir coulissant",
    "Caisson intérieur monté sur glissières — rangement accessible par l'extérieur",
    "98,5 × 41 × 31 cm",
    nb_planches=5, nb_carrelets=8, bloc_num=4, total_blocs=5
)

make_debit_slide("Tiroir coulissant", TIROIR_PLANCHES, TIROIR_CARRELETS_INT)
make_pieces_slide("Tiroir — planches", TIROIR_PLANCHES)
make_pieces_slide("Tiroir — carrelets",
    [(c[0], c[1], c[2], c[3], c[4], None) for c in TIROIR_CARRELETS_INT], is_carrelet=True)

make_steps_slide("Tiroir coulissant", [
    (1, "Ossature tiroir — carrelets [D6-D13]",
     "Assemble les 8 carrelets en cadre rectangulaire : 2 longueurs (D6, D7 — 98,5 cm) + 2 largeurs (D12, D13 — 38 cm) en bas, puis 4 montants verticaux (D8-D11 — 28 cm) aux angles. Vis 3,5×25 mm."),
    (2, "Poser les planches du caisson [D1-D5]",
     "Visse la planche bas [D1] sous l'ossature. Puis les 4 côtés : fond [D2], face [D3], droite [D4], gauche [D5]. Pré-perce systématiquement. 3 vis par planche."),
    (3, "Monter les glissières",
     "Fixe les glissières fixes (achetées) sur les carrelets rails B19-B22. Fixe les glissières coulissantes sur les côtés du tiroir. Teste le coulissement avant de tout serrer."),
    (4, "Insérer le tiroir",
     "Glisse le tiroir assemblé sur les rails. Il doit coulisser librement. Règle la butée si nécessaire."),
])


# ═══════════════════════════════════════
#  BLOC 5 — PEIGNE (LIT)
# ═══════════════════════════════════════
make_cover_slide(
    "Peigne — surface de couchage",
    "Le système lit peigne : 1 planche fixe + 2 peignes mobiles dépliables",
    "190 × 78 × 3 cm (déplié : 190 × 152 cm)",
    nb_planches=6, nb_carrelets=0, bloc_num=5, total_blocs=5
)

make_debit_slide("Peigne (lit)", PEIGNE)
make_pieces_slide("Peigne (lit)", PEIGNE)
make_steps_slide("Peigne (lit)", [
    (1, "Assembler les pieds + triangles [E5-E8]",
     "Colle et visse les triangles de renfort (E7, E8 — 10×10 cm) sur les pieds (E5, E6). Le triangle se fixe en haut du pied, côté intérieur. 2 vis 3,5×25 mm + colle à bois."),
    (2, "Fixer les pieds aux peignes mobiles [E3, E4]",
     "Visse chaque pied (E5 sur E3, E6 sur E4) perpendiculairement au peigne mobile. Le pied se fixe sur le bord long du peigne. 6 vis 3,5×30 mm par pied. Les dents du peigne doivent pointer vers le haut."),
    (3, "Poser la planche du dessus [E1]",
     "Place la planche du dessus lit [E1] (187 × 77,5 cm) sur le caisson. Elle repose sur les carrelets haut du bloc 2. Visse dans les carrelets — 6 vis 3,5×30 mm."),
    (4, "Poser le peigne fixe [E2]",
     "Le lit peigne fixée [E2] (190 × 74 cm) se pose PAR-DESSUS la planche du dessus, dents vers le haut. Visse sur les bords — 4 vis 3,5×30 mm. Les dents s'intercalent avec celles des peignes mobiles."),
    (5, "Installer les peignes mobiles",
     "Les 2 peignes mobiles (E3 droite, E4 gauche) viennent se replier sur le peigne fixe. Quand déplié, les dents du mobile s'emboîtent dans celles du fixe pour former une surface plane. Teste le dépliage avant de fixer définitivement les pieds."),
])


# ═══════════════════════════════════════
#  SLIDE FINALE — ACCESSOIRES + CHECKLIST
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE

add_text(slide, Cm(1.5), Cm(0.8), Cm(20), Cm(1.5),
         "Accessoires à acheter", size=18, bold=True, color=BLUE)

acc_items = [
    "• 2 paires de glissières coulissantes (100 cm)",
    "• 2 supports table extérieur",
    "• 1 table extérieur (90 × 40 cm) — ou planche découpée",
    "• 1 Boxio gaz + 1 Boxio robinet",
    "• Vis à bois 3,5 × 30 mm (×80) + 4 × 40 mm (×30) + 3,5 × 25 mm (×60)",
]
add_text(slide, Cm(1.5), Cm(2.5), Cm(15), Cm(6),
         "\n".join(acc_items), size=10, color=DARK)

# Checklist
add_rounded_rect(slide, Cm(17), Cm(2), Cm(15), Cm(10), fill_color=GREEN_BG)
add_text(slide, Cm(17.5), Cm(2.3), Cm(14), Cm(1),
         "Checklist finale ✅", size=14, bold=True, color=DARK)

checks = [
    "□  Équerrage vérifié (diagonales égales)",
    "□  Toutes les vis sont serrées",
    "□  Le tiroir coulisse sans forcer",
    "□  Les peignes mobiles se déploient correctement",
    "□  Les dents s'emboîtent sans jeu excessif",
    "□  La surface de couchage est plane",
    "□  Le meuble est stable et ne bouge pas",
    "□  Fixation au van vérifiée (équerres M6)",
]
add_text(slide, Cm(17.5), Cm(3.8), Cm(14), Cm(8),
         "\n".join(checks), size=10, color=DARK)

# Warning box
add_rounded_rect(slide, Cm(1.5), Cm(13), Cm(30), Cm(2.5), fill_color=WARN_BG)
add_text(slide, Cm(2), Cm(13.3), Cm(29), Cm(2),
         "⚠️  Fixation au van : Ce meuble doit être solidement fixé au plancher et aux parois avec des équerres métalliques et boulons M6. Ne jamais laisser un meuble libre en roulant.",
         size=10, bold=False, color=DARK)

add_text(slide, Cm(1.5), Cm(17.5), Cm(30), Cm(1),
         "Guide VBA Construction — Bloc Lit Peigne — vanzonexplorer.com",
         size=8, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════
prs.save(OUT_FILE)
print(f"✅ PPTX généré : {OUT_FILE}")
print(f"   Slides : {len(prs.slides)}")
print(f"   Taille : {os.path.getsize(OUT_FILE) / 1024:.0f} KB")
