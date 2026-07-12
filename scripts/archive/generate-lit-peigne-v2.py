#!/usr/bin/env python3
"""
Guide VBA Construction — Lit Peigne
Format: Portrait A4, style identique au guide_glaciere.pptx
Usage: DYLD_LIBRARY_PATH=/opt/homebrew/lib python3 scripts/generate-lit-peigne-v2.py
"""

import os, re
import cairosvg
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Config ───
SVG_DIR = "/Volumes/T7/pieces_van_svg2"
OUT_DIR = "/Users/julesgaveglio/vanzon-website-claude-code/public/docs"
PNG_CACHE = os.path.join(OUT_DIR, "png-cache-lit-v2")
OUT_FILE = os.path.join(OUT_DIR, "Guide-VBA-Lit-Peigne.pptx")
os.makedirs(PNG_CACHE, exist_ok=True)

# ─── A4 Portrait ───
PAGE_W = Cm(21.01)
PAGE_H = Cm(29.69)

BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x88, 0x88, 0x88)

# ─── Annotated SVG dir ───
ANNOTATED_DIR = os.path.join(OUT_DIR, "svg-annotated")

# ─── SVG → PNG ───
def svg_to_png(svg_filename, max_width=300):
    png_path = os.path.join(PNG_CACHE, svg_filename.replace(".svg", ".png"))
    if not os.path.exists(png_path):
        # Prefer annotated SVG if available
        annotated_path = os.path.join(ANNOTATED_DIR, svg_filename)
        if os.path.exists(annotated_path):
            svg_path = annotated_path
        else:
            svg_path = os.path.join(SVG_DIR, svg_filename)
        if os.path.exists(svg_path):
            cairosvg.svg2png(url=svg_path, write_to=png_path, output_width=max_width)
    return png_path if os.path.exists(png_path) else None


def get_svg_aspect(svg_filename):
    """Get width/height ratio from SVG (checks annotated first)."""
    annotated = os.path.join(ANNOTATED_DIR, svg_filename)
    if os.path.exists(annotated):
        return get_svg_aspect_from_file(annotated)
    return get_svg_aspect_from_file(os.path.join(SVG_DIR, svg_filename))

def get_svg_aspect_from_file(svg_path):
    """Get width/height ratio from SVG file path."""
    if not os.path.exists(svg_path):
        return 1.0
    with open(svg_path, 'r') as f:
        content = f.read(500)
    w_match = re.search(r'width="([\d.]+)', content)
    h_match = re.search(r'height="([\d.]+)', content)
    if w_match and h_match:
        w, h = float(w_match.group(1)), float(h_match.group(1))
        return w / h if h > 0 else 1.0
    return 1.0


# ─── Text helper ───
def add_text(slide, left, top, width, height, text, size=12, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    if size:
        p.font.size = Pt(size)
    if bold:
        p.font.bold = True
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return txBox


# ─── Piece data ───
# Each piece: (ref, name, dims_text, svg_file)
# Bloc 1: Base extérieure
BASE_PLANCHES = [
    ("A1", "Planche sol", "187 × 74,5 cm · ép. 1,5 cm", "066_Planche sol.svg"),
    ("A2", "Planche fond longueur", "187 × 36,5 cm · ép. 1,5 cm", "067_Planche fond longueur.svg"),
    ("A3", "Planche face longueur", "187 × 36,5 cm · ép. 1,5 cm", "070_Planche face longueur.svg"),
    ("A4", "Planche droite largeur", "77,5 × 38 cm · ép. 1,5 cm", "068_Planche droite largeur.svg"),
    ("A5", "Planche gauche largeur", "77,5 × 38 cm · ép. 1,5 cm", "069_Planche gauche largeur.svg"),
]

BASE_CARRELETS = [
    ("B1", "Longueur haut gauche", "187 cm", 1),
    ("B2", "Longueur face bas", "187 cm", 1),
    ("B3", "Longueur fond bas droit", "84,5 cm", 1),
    ("B4", "Longueur fond bas gauche", "16,5 cm", 1),
    ("B5", "Largeur milieu haut", "73 cm", 1),
    ("B6", "Largeur milieu haut 2", "73 cm", 1),
    ("B7", "Largeur bas droite", "71,5 cm", 1),
    ("B8", "Largeur 1", "71,5 cm", 1),
    ("B9", "Largeur milieu bas 2", "71,5 cm", 1),
    ("B10", "Droite largeur 2", "65,8 cm", 1),
    ("B11", "Vertical avant droite", "26,5 cm", 1),
    ("B12", "Vertical avant gauche", "26,5 cm", 1),
    ("B13", "Vertical fond droite", "32 cm", 1),
    ("B14", "Vertical intérieur fond", "32 cm", 1),
    ("B15", "Gauche vertical", "32 cm", 1),
    ("B16", "Vertical fond gauche petit", "32 cm", 1),
    ("B17", "Bas gauche largeur petit", "14,5 cm", 1),
    ("B18", "Petit gauche haut", "14,5 cm", 1),
    ("B19", "Long. gauche support tiroir", "101 cm", 1),
    ("B20", "Long. bas gauche support tiroir", "101 cm", 1),
    ("B21", "Long. droite support tiroir", "101 cm", 1),
    ("B22", "Long. bas droit support tiroir", "101 cm", 1),
    ("B23", "Vert. fond gauche support tiroir", "32 cm", 1),
    ("B24", "Vert. fond droit support tiroir", "33,5 cm", 1),
]

# Bloc 3: Cloisons intérieures
CLOISONS = [
    ("C1", "Support gauche glissière", "101 × 35 cm · ép. 1,5 cm", "093_Planche support gauche glissi_re.svg"),
    ("C2", "Planche droite intérieur", "101 × 35 cm · ép. 1,5 cm", "095_Planche droite int_rieur.svg"),
    ("C3", "Planche placard gauche", "74,5 × 35 cm · ép. 1,5 cm", "094_Planche placard gauche.svg"),
    ("C4", "Séparation boxio / rangement", "74,5 × 35 cm · ép. 1,5 cm", "096_Planche s_paration boxio _ rangement.svg"),
]

# Bloc 4: Tiroir coulissant
TIROIR_PLANCHES = [
    ("D1", "Planche bas tiroir", "98,5 × 41 cm · ép. 1,5 cm", "083_Planche bas cuisine.svg"),
    ("D2", "Planche fond tiroir", "44 × 31 cm · ép. 1,5 cm", "084_Planche fond cuisine.svg"),
    ("D3", "Planche face tiroir", "44 × 31 cm · ép. 1,5 cm", "087_Planche cuisine face.svg"),
    ("D4", "Planche droite tiroir", "98,5 × 31 cm · ép. 1,5 cm", "085_Planche droite cuisine.svg"),
    ("D5", "Planche gauche tiroir", "98,5 × 31 cm · ép. 1,5 cm", "086_Planche gauche cuisine.svg"),
]

TIROIR_CARRELETS = [
    ("D6", "Long. gauche tiroir", "98,5 cm", 1),
    ("D7", "Long. droit tiroir", "98,5 cm", 1),
    ("D8", "Fond gauche vertical", "28 cm", 1),
    ("D9", "Fond droit vertical", "28 cm", 1),
    ("D10", "Vertical avant gauche", "28 cm", 1),
    ("D11", "Vertical avant droit", "28 cm", 1),
    ("D12", "Fond largeur", "38 cm", 1),
    ("D13", "Largeur bas face", "38 cm", 1),
]

# Bloc 5: Peigne
PEIGNE_PLANCHES = [
    ("E1", "Planche du dessus lit", "187 × 77,5 cm · ép. 1,5 cm", "056_Planche du dessus lit.svg"),
    ("E2", "Lit peigne fixée", "190 × 74 cm · ép. 1,5 cm", "057_Lit peigne fixee.svg"),
    ("E3", "Peigne mobile droite", "95 × 74 cm · ép. 1,5 cm", "061_Peigne mobile partie droite.svg"),
    ("E4", "Peigne mobile gauche", "95 × 74 cm · ép. 1,5 cm", "062_Peigne mobile partie gauche.svg"),
    ("E5", "Pied lit droite", "95 × 38 cm · ép. 1,5 cm", "060_Pied lit peigne partie droite.svg"),
    ("E6", "Pied lit gauche", "95 × 38 cm · ép. 1,5 cm", "063_Pied lit peigne partie gauche.svg"),
    ("E7", "Triangle pied droit", "10 × 10 cm · ép. 1,5 cm", "059_Triangle pied droit.svg"),
    ("E8", "Triangle pied gauche", "10 × 10 cm · ép. 1,5 cm", "065_Triangle pied gauche.svg"),
]


# ═══════════════════════════════════════
#  BUILD PRESENTATION — Portrait A4
# ═══════════════════════════════════════
prs = Presentation()
prs.slide_width = PAGE_W
prs.slide_height = PAGE_H


def make_cover(title, subtitle_dims, matiere, nb_panneaux, nb_carrelets, carrelet_label="Carrelets 15 × 15 mm"):
    """Slide 1 style: centered title + info grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(slide, Cm(0), Cm(7.6), Cm(21), Cm(2), title, size=28, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Cm(0), Cm(9.7), Cm(21), Cm(1.3), "Guide de montage", size=16, align=PP_ALIGN.CENTER)

    # Info grid
    info = [
        ("Encombrement fini", subtitle_dims),
        ("Matière", matiere),
        ("Panneaux à découper", str(nb_panneaux)),
        (carrelet_label, str(nb_carrelets)),
    ]
    for i, (label, value) in enumerate(info):
        y = Cm(12.2 + i * 1.3)
        add_text(slide, Cm(4.1), y, Cm(6.6), Cm(1), label, size=11)
        add_text(slide, Cm(10.7), y, Cm(6.6), Cm(1), value, size=11)

    return slide


def make_debit(title, planches, carrelets=None):
    """Slide 2 style: tables for cutting list."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, Cm(1.5), Cm(1.3), Cm(17.8), Cm(1.5), title, size=20, bold=True)

    y = Cm(3.3)

    if planches:
        rows = len(planches) + 1
        tbl_shape = slide.shapes.add_table(rows, 4, Cm(1.5), y, Cm(17.9), Cm(0.6 * rows))
        tbl = tbl_shape.table
        tbl.columns[0].width = Cm(2)
        tbl.columns[1].width = Cm(7.5)
        tbl.columns[2].width = Cm(6)
        tbl.columns[3].width = Cm(2.4)

        for j, h in enumerate(["Rép.", "Panneau (CP 15 mm)", "Dimensions", "Qté"]):
            tbl.cell(0, j).text = h
        for i, p in enumerate(planches):
            tbl.cell(i+1, 0).text = p[0]
            tbl.cell(i+1, 1).text = p[1]
            tbl.cell(i+1, 2).text = p[2].split(" · ")[0]  # just dims, no thickness
            tbl.cell(i+1, 3).text = "1"

        y = y + Cm(0.6 * rows) + Cm(1)

    if carrelets:
        rows2 = len(carrelets) + 1
        # Check if fits on page, split if needed
        max_rows = min(rows2, int((Cm(28) - y) / Cm(0.5)))
        if max_rows < rows2:
            rows2 = max_rows

        tbl_shape2 = slide.shapes.add_table(rows2, 4, Cm(1.5), y, Cm(17.9), Cm(0.5 * rows2))
        tbl2 = tbl_shape2.table
        tbl2.columns[0].width = Cm(2)
        tbl2.columns[1].width = Cm(7.5)
        tbl2.columns[2].width = Cm(6)
        tbl2.columns[3].width = Cm(2.4)

        for j, h in enumerate(["Rép.", "Carrelet (15 × 15 mm)", "Longueur", "Qté"]):
            tbl2.cell(0, j).text = h
        for i, c in enumerate(carrelets[:rows2-1]):
            tbl2.cell(i+1, 0).text = c[0]
            tbl2.cell(i+1, 1).text = c[1]
            tbl2.cell(i+1, 2).text = c[2]
            tbl2.cell(i+1, 3).text = str(c[3])

        # If we couldn't fit all, make continuation slide
        remaining = carrelets[rows2-1:]
        if remaining:
            make_debit(f"{title} (suite)", None, remaining)

    return slide


# ─── Dimension annotations for complex pieces ───
# key: svg filename → list of (label, x_cm, y_cm) relative to image top-left
# These annotations are placed ON TOP of the piece image on detail slides
PIECE_ANNOTATIONS = {
    "066_Planche sol.svg": {
        "title": "Planche sol — avec encoche passage de roue",
        "notes": [
            ("1870", "longueur totale"),
            ("720,4", "largeur totale"),
            ("585", "largeur avant encoche"),
            ("135,4", "largeur encoche (720,4 − 585)"),
            ("835", "hauteur encoche (1705 − 870)"),
            ("870", "distance haut → début encoche"),
        ],
    },
    "067_Planche fond longueur.svg": {
        "title": "Planche fond — avec encoche passage de roue",
        "notes": [
            ("1870", "longueur totale"),
            ("365", "largeur totale"),
            ("44", "largeur partie réduite (zone milieu)"),
            ("321", "profondeur encoche (365 − 44)"),
            ("870", "distance haut → début encoche"),
            ("835", "hauteur encoche"),
        ],
    },
    "068_Planche droite largeur.svg": {
        "title": "Planche droite — avec chanfrein",
        "notes": [
            ("775", "largeur bas (totale)"),
            ("675", "largeur haut (avant chanfrein)"),
            ("380", "hauteur totale"),
            ("100 × 100", "chanfrein coin supérieur droit"),
        ],
    },
    "069_Planche gauche largeur.svg": {
        "title": "Planche gauche — avec chanfrein",
        "notes": [
            ("380", "largeur bas (totale)"),
            ("280", "largeur haut (avant chanfrein)"),
            ("775", "hauteur totale"),
            ("100 × 100", "chanfrein coin supérieur droit"),
        ],
    },
    "060_Pied lit peigne partie droite.svg": {
        "title": "Pied lit droite — avec encoche arrondie",
        "notes": [
            ("380", "largeur totale"),
            ("950", "hauteur totale"),
            ("300", "largeur avant encoche (en bas)"),
            ("50", "hauteur encoche arrondie"),
            ("R ≈ 40", "rayon arrondi de l'encoche"),
        ],
    },
    "063_Pied lit peigne partie gauche.svg": {
        "title": "Pied lit gauche — avec encoche arrondie",
        "notes": [
            ("380", "largeur totale"),
            ("950", "hauteur totale"),
            ("300", "largeur après encoche (en haut)"),
            ("50", "hauteur encoche arrondie"),
            ("R ≈ 40", "rayon arrondi de l'encoche"),
        ],
    },
    "061_Peigne mobile partie droite.svg": {
        "title": "Peigne mobile droite — forme peigne",
        "notes": [
            ("~950", "longueur planche brute"),
            ("~740", "largeur planche brute"),
            ("14 dents", "régulièrement espacées"),
            ("Découpe complexe", "tracer depuis le gabarit SVG"),
        ],
    },
    "062_Peigne mobile partie gauche.svg": {
        "title": "Peigne mobile gauche — forme peigne",
        "notes": [
            ("950", "longueur planche brute"),
            ("740", "largeur planche brute"),
            ("14 dents", "régulièrement espacées"),
            ("Découpe complexe", "tracer depuis le gabarit SVG"),
        ],
    },
    "057_Lit peigne fixee.svg": {
        "title": "Lit peigne fixée — forme peigne",
        "notes": [
            ("1900", "longueur totale"),
            ("740", "largeur totale"),
            ("28 dents", "s'emboîtent avec les peignes mobiles"),
            ("Découpe complexe", "tracer depuis le gabarit SVG"),
        ],
    },
}


def make_piece_detail(ref, name, dims, svg_file):
    """Full-page detail slide for a complex piece — annotated image fills the page."""
    annotations = PIECE_ANNOTATIONS.get(svg_file)
    if not annotations:
        return None

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_text(slide, Cm(1.5), Cm(0.8), Cm(17.8), Cm(1),
             f"{ref} — {annotations['title']}", size=14, bold=True)
    add_text(slide, Cm(1.5), Cm(1.8), Cm(17.8), Cm(0.6),
             f"{dims}  —  Toutes les cotes sont en mm", size=10, color=GREY)

    # Large annotated image (dimensions are ON the drawing)
    # Use annotated SVG → bigger render for detail
    annotated_png = os.path.join(PNG_CACHE, "detail_" + svg_file.replace(".svg", ".png"))
    annotated_svg = os.path.join(ANNOTATED_DIR, svg_file)
    if os.path.exists(annotated_svg) and not os.path.exists(annotated_png):
        cairosvg.svg2png(url=annotated_svg, write_to=annotated_png, output_width=1200)

    if os.path.exists(annotated_png):
        aspect = get_svg_aspect_from_file(annotated_svg)
        max_w = Cm(18)
        max_h = Cm(24)
        if aspect >= 1:
            img_w = max_w
            img_h = min(img_w / aspect, max_h)
        else:
            img_h = max_h
            img_w = min(img_h * aspect, max_w)
        # Center horizontally
        img_x = (PAGE_W - int(img_w)) // 2
        img_y = Cm(2.8)
        try:
            slide.shapes.add_picture(annotated_png, img_x, img_y, int(img_w), int(img_h))
        except Exception as e:
            print(f"  ⚠️ Detail image error {svg_file}: {e}")

    return slide


def make_pieces(title, pieces):
    """Slide 3 style: 3 columns × 2 rows of piece cards with PNG images."""
    items_per_slide = 6
    col_x = [Cm(1.5), Cm(7.7), Cm(14.0)]

    for page in range(0, len(pieces), items_per_slide):
        batch = pieces[page:page + items_per_slide]
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        page_num = page // items_per_slide + 1
        total = (len(pieces) + items_per_slide - 1) // items_per_slide
        suffix = f" ({page_num}/{total})" if total > 1 else ""
        add_text(slide, Cm(1.5), Cm(1.0), Cm(17.8), Cm(1.5), f"Pièces à découper — {title}{suffix}", size=18, bold=True)

        # Screenshot placeholder (first page only)
        if page == 0:
            from pptx.enum.shapes import MSO_SHAPE
            placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.5), Cm(2.8), Cm(8.1), Cm(4.6))
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
            placeholder.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            add_text(slide, Cm(1.5), Cm(4.2), Cm(8.1), Cm(1.5),
                     "Coller ici le screenshot\nde la vue éclatée", size=10, color=GREY, align=PP_ALIGN.CENTER)

        # Piece cards: row 1 starts at y=8.4, row 2 at y=15.2
        row_y = [Cm(8.4), Cm(15.2), Cm(22.0)]

        for i, piece in enumerate(batch):
            col = i % 3
            row = i // 3
            x = col_x[col]
            y_base = row_y[row]

            ref, name, dims, svg_file = piece[0], piece[1], piece[2], piece[3]

            # Ref label
            add_text(slide, x, y_base, Cm(1.3), Cm(0.6), ref, size=10, bold=True)
            # Name
            add_text(slide, x + Cm(1.2), y_base, Cm(4.8), Cm(0.6), name, size=10)

            # Image
            img_y = y_base + Cm(0.8)
            max_img_w = Cm(5)
            max_img_h = Cm(3.7)

            png_path = svg_to_png(svg_file, max_width=400)
            if png_path:
                aspect = get_svg_aspect(svg_file)
                if aspect >= 1:  # wider than tall
                    img_w = min(max_img_w, max_img_w)
                    img_h = min(img_w / aspect, max_img_h)
                else:  # taller than wide
                    img_h = min(max_img_h, max_img_h)
                    img_w = min(img_h * aspect, max_img_w)
                try:
                    slide.shapes.add_picture(png_path, x, img_y, int(img_w), int(img_h))
                except Exception as e:
                    print(f"  ⚠️ Image error {svg_file}: {e}")

            # Dimensions text
            dims_y = y_base + Cm(4.6)
            add_text(slide, x, dims_y, Cm(5.8), Cm(0.5), dims, size=8)

    # After grid slides, add detail slides for complex pieces
    for piece in pieces:
        ref, name, dims, svg_file = piece[0], piece[1], piece[2], piece[3]
        if svg_file in PIECE_ANNOTATIONS:
            make_piece_detail(ref, name, dims, svg_file)


def make_steps(title, steps):
    """Slide 4 style: numbered steps with title + description."""
    steps_per_slide = 4

    for page in range(0, len(steps), steps_per_slide):
        batch = steps[page:page + steps_per_slide]
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        page_num = page // steps_per_slide + 1
        total = (len(steps) + steps_per_slide - 1) // steps_per_slide
        suffix = f" ({page_num}/{total})" if total > 1 else ""
        add_text(slide, Cm(1.5), Cm(1.3), Cm(17.8), Cm(1.5), f"Étapes de montage — {title}{suffix}", size=20, bold=True)

        for i, (num, step_title, instruction) in enumerate(batch):
            y = Cm(3.6 + i * 3.3)
            # Number
            add_text(slide, Cm(1.5), y, Cm(1.5), Cm(1.5), str(num), size=18, bold=True, align=PP_ALIGN.CENTER)
            # Title
            add_text(slide, Cm(3.3), y, Cm(15.2), Cm(1), step_title, size=13, bold=True)
            # Instruction
            add_text(slide, Cm(3.3), y + Cm(1), Cm(16), Cm(2), instruction, size=10)


# ═══════════════════════════════════════
#  BLOC 1 — BASE EXTÉRIEURE
# ═══════════════════════════════════════
make_cover("Bloc Lit Peigne\nBase extérieure", "187 × 77,5 × 38 cm",
           "CP peuplier 15 mm", nb_panneaux=5, nb_carrelets=24)
make_debit("Liste de débit — Base extérieure", BASE_PLANCHES, BASE_CARRELETS)
make_pieces("Base extérieure", BASE_PLANCHES)
make_steps("Base extérieure", [
    (1, "Poser la planche sol [A1]",
     "Place la planche sol à plat. Attention : encoche pour le passage de roue — vérifie l'orientation."),
    (2, "Monter les côtés [A4] et [A5]",
     "Visse les planches droite et gauche verticalement sur les bords de A1. Pré-perce à 3 mm, 4 vis 4×40 mm par côté."),
    (3, "Fixer le fond [A2]",
     "Glisse la planche fond longueur à l'arrière. Visse dans les chants des côtés — 4 vis 3,5×30 mm par côté."),
    (4, "Fermer avec la face [A3]",
     "Fixe la planche face longueur à l'avant. Vérifie l'équerrage (diagonales) avant de serrer."),
])


# ═══════════════════════════════════════
#  BLOC 2 — OSSATURE CARRELETS (déjà dans la liste de débit du bloc 1)
#  → slide étapes séparée
# ═══════════════════════════════════════
make_steps("Ossature carrelets", [
    (5, "Cadre bas [B2, B3, B4, B7-B9]",
     "Visse les carrelets horizontaux au fond de la base. B2 (face, 187 cm), B3+B4 (fond), B7-B9 (largeurs). 2 vis 3,5×25 mm par carrelet."),
    (6, "Montants verticaux [B11-B16]",
     "Fixe les 6 carrelets verticaux dans les angles et jonctions internes. Visse par le bas + par le côté."),
    (7, "Cadre haut [B1, B5, B6, B10]",
     "Pose les carrelets du haut sur les montants. B1 (187 cm), B5+B6 (73 cm), B10 (65,8 cm). 2 vis par jonction."),
    (8, "Petits carrelets [B17, B18]",
     "Fixe les 2 petits carrelets (14,5 cm) en bas à gauche (B17) et en haut à gauche (B18)."),
])
make_steps("Rails support tiroir", [
    (9, "Rails tiroir [B19-B24]",
     "Visse les 4 carrelets horizontaux (B19-B22, ~101 cm) parallèles — ils forment les rails du tiroir. Puis les 2 verticaux (B23-B24) à l'arrière. Le tiroir doit coulisser sans forcer."),
])


# ═══════════════════════════════════════
#  BLOC 3 — CLOISONS INTÉRIEURES
# ═══════════════════════════════════════
make_cover("Bloc Lit Peigne\nCloisons intérieures", "—",
           "CP peuplier 15 mm", nb_panneaux=4, nb_carrelets=0)
make_debit("Liste de débit — Cloisons", CLOISONS)
make_pieces("Cloisons", CLOISONS)
make_steps("Cloisons intérieures", [
    (1, "Cloisons longitudinales [C1] et [C2]",
     "Visse C1 (support glissière) et C2 (droite intérieur) verticalement sur les carrelets. 4 vis 3,5×30 mm par planche."),
    (2, "Cloisons transversales [C3] et [C4]",
     "Fixe C3 (placard gauche) et C4 (séparation boxio) perpendiculairement. Visse dans les carrelets + cloisons longitudinales."),
])


# ═══════════════════════════════════════
#  BLOC 4 — TIROIR COULISSANT
# ═══════════════════════════════════════
make_cover("Bloc Lit Peigne\nTiroir coulissant", "98,5 × 41 × 31 cm",
           "CP peuplier 15 mm", nb_panneaux=5, nb_carrelets=8)
make_debit("Liste de débit — Tiroir", TIROIR_PLANCHES, TIROIR_CARRELETS)
make_pieces("Tiroir — planches", TIROIR_PLANCHES)
make_steps("Tiroir coulissant", [
    (1, "Ossature tiroir [D6-D13]",
     "Assemble les 8 carrelets en cadre : 2 longueurs (D6-D7, 98,5 cm) + 2 largeurs (D12-D13, 38 cm) + 4 verticaux (D8-D11, 28 cm)."),
    (2, "Planches du caisson [D1-D5]",
     "Visse la planche bas [D1], puis les 4 côtés : fond [D2], face [D3], droite [D4], gauche [D5]. 3 vis par planche."),
    (3, "Monter les glissières",
     "Fixe les glissières fixes sur les rails B19-B22. Fixe les coulissantes sur le tiroir. Teste le coulissement."),
    (4, "Insérer le tiroir",
     "Glisse le tiroir sur les rails. Il doit coulisser librement."),
])


# ═══════════════════════════════════════
#  BLOC 5 — PEIGNE (LIT)
# ═══════════════════════════════════════
make_cover("Bloc Lit Peigne\nPeigne — surface de couchage", "190 × 152 cm (déplié)",
           "CP peuplier 15 mm", nb_panneaux=8, nb_carrelets=0)
make_debit("Liste de débit — Peigne", PEIGNE_PLANCHES)
make_pieces("Peigne", PEIGNE_PLANCHES)
make_steps("Peigne (lit)", [
    (1, "Pieds + triangles [E5-E8]",
     "Colle et visse les triangles de renfort (E7, E8) sur les pieds (E5, E6). 2 vis 3,5×25 mm + colle à bois."),
    (2, "Pieds sur peignes mobiles [E3, E4]",
     "Visse chaque pied perpendiculairement au peigne mobile. 6 vis 3,5×30 mm par pied. Dents vers le haut."),
    (3, "Planche du dessus [E1]",
     "Place E1 (187 × 77,5 cm) sur le caisson, sur les carrelets haut. Visse — 6 vis 3,5×30 mm."),
    (4, "Peigne fixe [E2]",
     "Pose E2 (190 × 74 cm) par-dessus E1, dents vers le haut. Visse sur les bords — 4 vis 3,5×30 mm."),
    (5, "Installer les peignes mobiles",
     "Les peignes E3 et E4 se replient sur le fixe. Dépliés, les dents s'emboîtent pour former une surface plane. Teste le dépliage."),
])


# ═══════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════
prs.save(OUT_FILE)
print(f"✅ PPTX généré : {OUT_FILE}")
print(f"   Format : Portrait A4 ({prs.slide_width/360000:.0f} × {prs.slide_height/360000:.0f} cm)")
print(f"   Slides : {len(prs.slides)}")
print(f"   Taille : {os.path.getsize(OUT_FILE) / 1024:.0f} KB")
